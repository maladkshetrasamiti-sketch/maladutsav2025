#!/usr/bin/env python3
"""
student_ppt.py

Generate PowerPoint presentation with student details and photos:
- Read student data from student_ppt_csv.csv
- Load photos from Final_Processed_photos folder
- Create one slide per student with details and photo
- Save as PowerPoint presentation
"""

import os
import datetime
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import subprocess
import shutil
import time
import sys

# Configuration
CSV_FILE = "student_ppt_csv.csv"
PHOTOS_FOLDER = "Final_Processed_photos"
OUTPUT_FILE = "Students_Presentation.pptx"
MPM_LOGO = "MPM_logo.png"  # place this file next to the script or inside Final_Processed_photos

# Slide layout settings
TITLE_FONT_SIZE = Pt(32)
DETAIL_FONT_SIZE = Pt(18)
PHOTO_WIDTH = Inches(2.88)   # 0.8 times previous width
PHOTO_HEIGHT = Inches(3.84)  # 0.8 times previous height

def read_student_data():
    """Read and validate student data from CSV file."""
    encodings = ['utf-8', 'latin-1', 'cp1252']
    
    for encoding in encodings:
        try:
            df = pd.read_csv(CSV_FILE, encoding=encoding)
            # Rename 'Student Name' column to 'Name' if it exists
            if 'Student Name' in df.columns:
                df = df.rename(columns={'Student Name': 'Name'})
            
            if 'Name' not in df.columns:
                print("⚠️ Error: CSV file must contain a 'Student Name' or 'Name' column")
                return None
                
            print(f"✓ Read {len(df)} student records from {CSV_FILE} using {encoding} encoding")
            return df
        except UnicodeDecodeError:
            continue
        except Exception as e:
            print(f"⚠️ Error reading {CSV_FILE}: {str(e)}")
            return None
    
    print(f"⚠️ Could not read {CSV_FILE} with any of the attempted encodings")
    return None

def get_photo_path(student_name):
    """Get the photo file path for a student."""
    # Look for common image extensions
    for ext in ['.jpg', '.jpeg', '.png']:
        photo_path = os.path.join(PHOTOS_FOLDER, f"{student_name}{ext}")
        if os.path.exists(photo_path):
            return photo_path
    return None


def find_logo():
    """Try to locate the MPM logo file in a few likely places.

    Returns the path to the logo file or None if not found.
    """
    candidates = [MPM_LOGO, os.path.join(PHOTOS_FOLDER, MPM_LOGO), os.path.join(os.getcwd(), MPM_LOGO)]
    for c in candidates:
        if os.path.exists(c):
            return c
    return None

def create_student_slide(prs, student):
    """Create a slide for a single student with photo and details."""
    # Skip if student has no name
    if pd.isna(student.get('Name')):
        print("⚠️ Skipping student with missing name")
        return
        
    # Use a blank slide layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Set slide background to black
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    # Add header box - positioned to start after the logo
    header_box = slide.shapes.add_textbox(
        Inches(2.2), Inches(0.2),  # Moved right to accommodate logo (1.8" logo + 0.4" spacing)
        Inches(7.3), Inches(1.8)  # Width adjusted to maintain right margin
    )
    header_frame = header_box.text_frame
    header_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add header lines
    header_lines = [
        "MAHESHWARI PRAGATI MANDAL",
        "MALAD SAMITIES",
        "MALAD UTSAV 2025"
    ]
    
    for i, line in enumerate(header_lines):
        if i == 0:
            p = header_frame.paragraphs[0]
        else:
            p = header_frame.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        # Increase 'MALAD UTSAV 2025' by 4pt (was 28 -> now 32)
        run.font.size = Pt(32) if (i == 0 or i == 2) else Pt(28)
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True
    
    # Add logo at top-left if available
    logo_path = find_logo()
    if logo_path:
        try:
            # 1.5 times the previous size (was 1.2")
            slide.shapes.add_picture(logo_path, Inches(0.2), Inches(0.2), height=Inches(1.8))
        except Exception as e:
            print(f"⚠️ Warning: could not add logo: {e}")

    # Calculate vertical center alignment with photo
    photo_top = Inches(2)  # Photo's top position
    photo_height = PHOTO_HEIGHT
    photo_center = photo_top + (photo_height / 2)
    details_height = Inches(4)  # Height of details box
    details_top = photo_center - (details_height / 2)  # Center details with photo

    # First, create an invisible textbox to measure text height
    temp_box = slide.shapes.add_textbox(
        Inches(0), Inches(0), Inches(5), Inches(4)
    )
    temp_frame = temp_box.text_frame
    temp_frame.word_wrap = True

    # Add temporary text to measure height
    temp_p = temp_frame.paragraphs[0]
    temp_p.text = "Name"
    temp_p.font.size = DETAIL_FONT_SIZE
    
    # Add 4 paragraphs to measure total height
    for _ in range(3):
        p = temp_frame.add_paragraph()
        p.text = "Sample"
        p.font.size = DETAIL_FONT_SIZE
    
    # Get text height (approximate)
    text_height = Inches(1.2)  # Height of 4 lines with spacing
    
    # Calculate center position relative to photo
    photo_top = Inches(2)
    photo_center = photo_top + (PHOTO_HEIGHT / 2)
    text_top = photo_center - (text_height / 2)
    
    # Remove temporary textbox
    slide.shapes._spTree.remove(temp_box._element)
    
    # Add details textbox at calculated position with adjusted width and position
    details_box = slide.shapes.add_textbox(
        Inches(0.5), text_top,  # Vertically centered with photo
        Inches(5.5), text_height  # Width increased due to smaller photo and repositioning
    )
    details_frame = details_box.text_frame
    details_frame.word_wrap = True
    details_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertical center alignment

    # Add student name with stars above and below (instead of inline)
    # Top stars
    top_stars = details_frame.paragraphs[0]
    top_stars.text = "⭐⭐⭐⭐⭐"
    top_stars.alignment = PP_ALIGN.CENTER
    top_stars.font.size = Pt(14)
    top_stars.font.bold = True
    top_stars.font.color.rgb = RGBColor(255, 0, 0)
    top_stars.space_after = Pt(6)

    # Name line
    name_para = details_frame.add_paragraph()
    name_para.text = student['Name'].upper()
    name_para.alignment = PP_ALIGN.CENTER
    name_para.font.size = Pt(28)
    name_para.font.bold = True
    name_para.font.color.rgb = RGBColor(255, 0, 0)
    name_para.space_after = Pt(6)

    # Bottom stars
    bottom_stars = details_frame.add_paragraph()
    bottom_stars.text = "⭐⭐⭐⭐⭐"
    bottom_stars.alignment = PP_ALIGN.CENTER
    bottom_stars.font.size = Pt(14)
    bottom_stars.font.bold = True
    bottom_stars.font.color.rgb = RGBColor(255, 0, 0)
    bottom_stars.space_after = Pt(18)
    
    # Map field names to desired display names
    field_display_names = {
        'Parent Name': 'Parent Name',
        'School Name': 'School Name',
        'Grade Standard': 'Grade Standard'
    }

    # Add each detail as a new paragraph
    for field, display_name in field_display_names.items():
        if field in student.index and not pd.isna(student[field]):
            p = details_frame.add_paragraph()
            p.text = f"{display_name}: {student[field]}"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = DETAIL_FONT_SIZE
            p.font.color.rgb = RGBColor(255, 255, 255)  # White text
            p.space_after = Pt(12)  # Consistent spacing between lines
    
    # Add photo on the right if available
    photo_path = get_photo_path(student['Name'])
    if photo_path:
        left = Inches(6.5)  # Adjusted for smaller photo to maintain right alignment
        top = Inches(2.2)   # Adjusted to maintain vertical center with smaller height
        try:
            # Add a thin border around the photo by adding a rectangle first
            border = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left - Inches(0.01),  # Slightly larger than photo
                top - Inches(0.01),
                PHOTO_WIDTH + Inches(0.02),
                PHOTO_HEIGHT + Inches(0.02)
            )
            border.fill.background()  # No fill
            border.line.color.rgb = RGBColor(255, 255, 255)  # White border
            border.line.width = Pt(2)  # Border thickness
            
            # Add the photo on top
            slide.shapes.add_picture(
                photo_path, 
                left, top,
                width=PHOTO_WIDTH,
                height=PHOTO_HEIGHT
            )
        except Exception as e:
            print(f"⚠️ Error adding photo for {student['Name']}: {str(e)}")
    else:
        print(f"⚠️ No photo found for {student['Name']}")

def save_ppt_as_pdf(pptx_path):
    """Attempt to convert a .pptx to .pdf.

    Strategies tried (in order):
    1. comtypes (PowerPoint COM automation)
    2. win32com (pywin32)
    3. LibreOffice/soffice command-line

    Returns the path to the PDF if successful, otherwise None.
    """
    pptx_path = os.path.abspath(pptx_path)
    pdf_path = os.path.splitext(pptx_path)[0] + '.pdf'

    # Try comtypes (works without requiring full pywin32 in some setups)
    try:
        import comtypes.client
        print("→ Converting PPTX to PDF using PowerPoint (comtypes)...")
        powerpoint = comtypes.client.CreateObject('Powerpoint.Application')
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
        # 32 is the enum for PDF format
        presentation.SaveAs(pdf_path, 32)
        presentation.Close()
        powerpoint.Quit()
        time.sleep(0.5)
        if os.path.exists(pdf_path):
            return pdf_path
    except Exception:
        pass

    # Try pywin32 (win32com)
    try:
        from win32com import client
        print("→ Converting PPTX to PDF using PowerPoint (win32com)...")
        powerpoint = client.Dispatch('Powerpoint.Application')
        powerpoint.Visible = True
        presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
        presentation.SaveAs(pdf_path, 32)
        presentation.Close()
        powerpoint.Quit()
        time.sleep(0.5)
        if os.path.exists(pdf_path):
            return pdf_path
    except Exception:
        pass

    # Try LibreOffice / soffice
    soffice = shutil.which('soffice') or shutil.which('libreoffice')
    if soffice:
        try:
            print(f"→ Converting PPTX to PDF using LibreOffice at: {soffice}")
            cmd = [soffice, '--headless', '--convert-to', 'pdf', '--outdir', os.path.dirname(pptx_path), pptx_path]
            subprocess.run(cmd, check=True)
            time.sleep(0.5)
            if os.path.exists(pdf_path):
                return pdf_path
        except Exception:
            pass

    # If all methods fail, return None
    print("⚠️ Could not convert PPTX to PDF automatically. Install 'comtypes' or 'pywin32' or LibreOffice to enable automatic conversion.")
    return None

def create_presentation():
    """Create the full presentation with all student slides."""
    # Read student data
    df = read_student_data()
    if df is None:
        return False
    
    # Create new presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # (No opening title slide) Slides will start directly with student entries.
    
    # Create a slide for each student
    success_count = 0
    for _, student in df.iterrows():
        try:
            create_student_slide(prs, student)
            success_count += 1
        except Exception as e:
            print(f"⚠️ Error creating slide for {student.get('Name', 'Unknown')}: {str(e)}")
    
    # Save the presentation using a timestamped filename to avoid overwrite/permission issues
    try:
        base, ext = os.path.splitext(OUTPUT_FILE)
        timestamp = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
        save_name = f"{base}_{timestamp}{ext}"
        prs.save(save_name)
        print(f"\n✓ Successfully created {success_count} slides")
        print(f"✓ Presentation saved as: {save_name}")
        # Attempt to save as PDF as well
        try:
            # Convert PPTX to PDF (default output name will mirror the PPTX base)
            pdf_path = save_ppt_as_pdf(save_name)
            # Change PDF base name from e.g. 'Students_Presentation_...' to 'Students_PDF_...'
            desired_pdf_base = base.replace('Presentation', 'PDF') if 'Presentation' in base else base + '_PDF'
            desired_pdf_name = f"{desired_pdf_base}_{timestamp}.pdf"
            desired_pdf_path = os.path.join(os.path.dirname(save_name), desired_pdf_name)

            if pdf_path and os.path.exists(pdf_path):
                # Rename/move to desired name if needed
                if os.path.abspath(pdf_path) != os.path.abspath(desired_pdf_path):
                    try:
                        # Overwrite if exists
                        if os.path.exists(desired_pdf_path):
                            os.remove(desired_pdf_path)
                        os.replace(pdf_path, desired_pdf_path)
                        pdf_path = desired_pdf_path
                    except Exception:
                        # If rename fails, keep original
                        pass
                print(f"✓ PDF saved as: {pdf_path}")
            else:
                print("⚠️ PDF conversion was not completed automatically.")
        except Exception as e:
            print(f"⚠️ Exception while converting to PDF: {e}")
        return True
    except Exception as e:
        print(f"⚠️ Error saving presentation: {str(e)}")
        # As a last resort try saving to temp directory
        try:
            import tempfile
            temp_path = os.path.join(tempfile.gettempdir(), os.path.basename(save_name))
            prs.save(temp_path)
            print(f"✓ Saved presentation to temp path: {temp_path}")
            return True
        except Exception as e2:
            print(f"⚠️ Failed to save presentation to temp path: {e2}")
            print("→ Please ensure you have write permission to the folder or close any open file with the same name and re-run the script.")
            return False

def main():
    """Main entry point."""
    # Verify folders and files exist
    if not os.path.exists(CSV_FILE):
        print(f"⚠️ CSV file not found: {CSV_FILE}")
        return
    
    if not os.path.exists(PHOTOS_FOLDER):
        print(f"⚠️ Photos folder not found: {PHOTOS_FOLDER}")
        return
    
    # Create presentation
    create_presentation()

if __name__ == '__main__':
    main()